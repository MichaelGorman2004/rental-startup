#!/usr/bin/env python3
"""
VenueLink Pitch Deck Generator
Creates a professional PowerPoint presentation for a 2-minute video pitch.
6 slides total: Intro + 4 content slides + Conclusion
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_venuelink_pitch():
    """Generate the complete VenueLink pitch deck."""

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Define color scheme (matching mockup aesthetic)
    DARK_BG = RGBColor(20, 20, 20)  # #141414
    DARKER_BG = RGBColor(10, 10, 10)  # #0a0a0a
    WHITE_TEXT = RGBColor(255, 255, 255)  # #FFFFFF
    GRAY_TEXT = RGBColor(136, 136, 136)  # #888888
    ACCENT_GRAY = RGBColor(42, 42, 42)  # #2a2a2a

    # Slide 1: Title/Intro
    slide1 = create_title_slide(prs, DARK_BG, WHITE_TEXT, GRAY_TEXT)

    # Slide 2: The Problem (+ for whom)
    slide2 = create_problem_slide(prs, DARK_BG, WHITE_TEXT, GRAY_TEXT, ACCENT_GRAY)

    # Slide 3: The Solution (with screenshots)
    slide3 = create_solution_slide(prs, DARK_BG, WHITE_TEXT, GRAY_TEXT, ACCENT_GRAY)

    # Slide 4: Current Progress
    slide4 = create_progress_slide(prs, DARK_BG, WHITE_TEXT, GRAY_TEXT, ACCENT_GRAY)

    # Slide 5: Why We're Building This
    slide5 = create_why_slide(prs, DARK_BG, WHITE_TEXT, GRAY_TEXT, ACCENT_GRAY)

    # Slide 6: Conclusion
    slide6 = create_closing_slide(prs, DARK_BG, WHITE_TEXT, GRAY_TEXT, ACCENT_GRAY)

    # Save presentation
    output_path = 'VenueLink_Pitch.pptx'
    prs.save(output_path)
    print(f"✓ PowerPoint presentation created: {output_path}")
    print(f"✓ Total slides: {len(prs.slides)}")
    print(f"✓ Estimated presentation time: ~2 minutes (~20 seconds per slide)")

    return output_path


def create_title_slide(prs, bg_color, white, gray):
    """Slide 1: Title/Intro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Title: VenueLink
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "VenueLink"
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = white

    # Tagline
    tagline_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.4), Inches(8), Inches(0.4)
    )
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Connecting Student Organizations with Local Venues"
    p = tagline_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.color.rgb = gray

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.0), Inches(8), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Marketplace for Event Venue Rentals"
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = gray

    # Team
    team_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(8), Inches(0.4)
    )
    team_frame = team_box.text_frame
    team_frame.text = "Michael Gorman & Larry Stelzer"
    p = team_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.color.rgb = white

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """Hi, I'm Michael, and this is VenueLink. We're building a marketplace that connects student organizations with local venues for event bookings. We're making it fast and easy for student groups to find and book bars, restaurants, and event spaces."""

    return slide


def create_problem_slide(prs, bg_color, white, gray, accent):
    """Slide 2: The Problem (+ for whom)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Header
    header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    header_frame = header_box.text_frame
    header_frame.text = "The Problem: Coordination Chaos"
    p = header_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = white

    # Personal context
    context_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.0), Inches(8), Inches(0.5)
    )
    cf = context_box.text_frame
    cf.text = "As fraternity brothers, we've experienced this problem firsthand"
    p = cf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = gray

    # Left Column - Student Organizations
    left_x = 0.8
    left_y = 1.8

    # Column header
    left_header = slide.shapes.add_textbox(
        Inches(left_x), Inches(left_y), Inches(4), Inches(0.4)
    )
    lh_frame = left_header.text_frame
    lh_frame.text = "Student Organizations"
    p = lh_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = white

    # Problems list
    problems_student = [
        "Hours spent calling/emailing venues",
        "No way to compare options easily",
        "Unclear pricing and availability",
        "Disorganized booking tracking",
        "Last-minute scrambling for events"
    ]

    left_box = slide.shapes.add_textbox(
        Inches(left_x), Inches(left_y + 0.6), Inches(4), Inches(3.2)
    )
    tf = left_box.text_frame
    tf.word_wrap = True

    for i, problem in enumerate(problems_student):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = f"• {problem}"
        p.font.size = Pt(15)
        p.font.color.rgb = white
        p.space_after = Pt(10)

    # Right Column - Venues
    right_x = 5.2
    right_y = 1.8

    # Column header
    right_header = slide.shapes.add_textbox(
        Inches(right_x), Inches(right_y), Inches(4), Inches(0.4)
    )
    rh_frame = right_header.text_frame
    rh_frame.text = "Local Venues"
    p = rh_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = white

    # Problems list
    problems_venue = [
        "Time-consuming coordination",
        "No formal booking system",
        "Payment and no-show issues",
        "Difficult to reach student groups",
        "Missed revenue opportunities"
    ]

    right_box = slide.shapes.add_textbox(
        Inches(right_x), Inches(right_y + 0.6), Inches(4), Inches(3.2)
    )
    tf = right_box.text_frame
    tf.word_wrap = True

    for i, problem in enumerate(problems_venue):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = f"• {problem}"
        p.font.size = Pt(15)
        p.font.color.rgb = white
        p.space_after = Pt(10)

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """We know this problem intimately because we've lived it. As fraternity brothers, we've spent countless hours calling bar after bar trying to find one with availability for our events. Student organizations waste time on email and phone tag with no centralized way to compare venues or check availability. Meanwhile, venues struggle with disorganized communication, no formal booking process, and difficulty reaching student groups. There's no platform connecting these two sides."""

    return slide


def create_solution_slide(prs, bg_color, white, gray, accent):
    """Slide 3: The Solution (with screenshot placeholders)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Header
    header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
    )
    header_frame = header_box.text_frame
    header_frame.text = "The Solution: VenueLink Marketplace"
    p = header_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = white

    # How it works - 3 steps (compact)
    steps_y = 0.95
    steps = [
        ("1. Discover", "Browse venues with filters"),
        ("2. Request", "Submit booking details"),
        ("3. Confirm", "Pay and get confirmed")
    ]

    step_width = 2.8
    gap = 0.3
    start_x = 1.0

    for i, (title, desc) in enumerate(steps):
        x_pos = start_x + (i * (step_width + gap))

        # Step box
        step_box = slide.shapes.add_textbox(
            Inches(x_pos), Inches(steps_y), Inches(step_width), Inches(0.5)
        )
        sf = step_box.text_frame
        sf.text = title
        p = sf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = white

        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(x_pos), Inches(steps_y + 0.4), Inches(step_width), Inches(0.3)
        )
        df = desc_box.text_frame
        df.text = desc
        df.word_wrap = True
        p = df.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.color.rgb = gray

    # Screenshots section
    screenshot_y = 2.0
    screenshot_label = slide.shapes.add_textbox(
        Inches(0.5), Inches(screenshot_y), Inches(9), Inches(0.3)
    )
    sl_frame = screenshot_label.text_frame
    sl_frame.text = "What We've Built:"
    p = sl_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = white

    # Screenshot placeholders (2-3 screenshots)
    screenshots = [
        "Login Page",
        "Student Org Dashboard"
    ]

    screenshot_width = 4.2
    screenshot_height = 2.4
    screenshot_gap = 0.6
    screenshot_start_x = 1.0
    screenshot_start_y = 2.5

    for i, label in enumerate(screenshots):
        x_pos = screenshot_start_x + (i * (screenshot_width + screenshot_gap))

        # Placeholder box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(screenshot_start_y),
            Inches(screenshot_width), Inches(screenshot_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(26, 26, 26)
        box.line.color.rgb = accent

        # Label
        text_frame = box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.text = f"[{label} Screenshot]"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = gray

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """VenueLink is a centralized marketplace that solves this problem. Student organizations can browse all available venues in one place, filter by their needs like capacity and price, submit booking requests with their event details, and once a venue accepts, pay a booking fee to confirm. It's like Airbnb but designed specifically for student organizations and local venues. We've already built a working authentication system and student dashboard that you can see here."""

    return slide


def create_progress_slide(prs, bg_color, white, gray, accent):
    """Slide 4: Current Progress"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Header
    header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    header_frame = header_box.text_frame
    header_frame.text = "Current Progress"
    p = header_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = white

    # What's Built section
    built_y = 1.1
    built_header = slide.shapes.add_textbox(
        Inches(0.8), Inches(built_y), Inches(4.2), Inches(0.4)
    )
    bh_frame = built_header.text_frame
    bh_frame.text = "✅ What We've Built"
    p = bh_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = white

    built_items = [
        "Authentication system with .edu verification",
        "Student organization dashboard",
        "Backend API for venue management",
        "Database schema (4 tables, PostgreSQL)",
        "AI-powered development workflow"
    ]

    built_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(built_y + 0.5), Inches(4.2), Inches(2.2)
    )
    btf = built_box.text_frame
    btf.word_wrap = True

    for i, item in enumerate(built_items):
        if i > 0:
            btf.add_paragraph()
        p = btf.paragraphs[i]
        p.text = f"• {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = white
        p.space_after = Pt(8)

    # What's Left section
    left_y = 1.1
    left_header = slide.shapes.add_textbox(
        Inches(5.2), Inches(left_y), Inches(4.0), Inches(0.4)
    )
    lh_frame = left_header.text_frame
    lh_frame.text = "🔨 What's Left to Build"
    p = lh_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = white

    left_items = [
        "Venue discovery page with search/filters",
        "Booking request flow",
        "Venue admin dashboard",
        "Payment processing integration"
    ]

    left_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(left_y + 0.5), Inches(4.0), Inches(1.5)
    )
    ltf = left_box.text_frame
    ltf.word_wrap = True

    for i, item in enumerate(left_items):
        if i > 0:
            ltf.add_paragraph()
        p = ltf.paragraphs[i]
        p.text = f"• {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = gray
        p.space_after = Pt(8)

    # Validation section
    validation_y = 3.8
    validation_header = slide.shapes.add_textbox(
        Inches(5.2), Inches(validation_y), Inches(4.0), Inches(0.4)
    )
    vh_frame = validation_header.text_frame
    vh_frame.text = "📋 Validation Plan"
    p = vh_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = white

    validation_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(validation_y + 0.5), Inches(4.0), Inches(0.8)
    )
    vf = validation_box.text_frame
    vf.word_wrap = True
    vf.text = "• Survey prepared for student orgs & venues\n• Planning interviews with 10+ bars and 10+ student organizations at Ohio State"
    p = vf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = white

    # Tech stack
    tech_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.9), Inches(8.5), Inches(0.3)
    )
    tech_frame = tech_box.text_frame
    tech_frame.text = "Tech: React + TypeScript + Mantine UI • FastAPI + PostgreSQL • Railway"
    p = tech_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = gray

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """Here's what we've built so far. We have a fully functional authentication system that verifies .edu emails, a working student organization dashboard, a complete backend API for venue management, and our database schema. We're using modern technologies like React, TypeScript, FastAPI, and PostgreSQL with an AI-powered development workflow. What's left to build includes the venue discovery page, the full booking request flow, and the venue admin dashboard. For validation, we've prepared comprehensive surveys and are planning to interview over 10 bars and 10 student organizations at Ohio State to validate our assumptions before full launch."""

    return slide


def create_why_slide(prs, bg_color, white, gray, accent):
    """Slide 5: Why We're Building This"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Header
    header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    header_frame = header_box.text_frame
    header_frame.text = "Why We're Building This"
    p = header_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = white

    # Personal Connection
    personal_y = 1.2
    personal_header = slide.shapes.add_textbox(
        Inches(1.5), Inches(personal_y), Inches(7), Inches(0.4)
    )
    ph_frame = personal_header.text_frame
    ph_frame.text = "👥 Personal Connection"
    p = ph_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = white

    personal_items = [
        "We're fraternity brothers who've lived this exact problem",
        "Spent hours calling 10+ bars just to find one available date",
        "Experienced the frustration of coordinating events firsthand",
        "We ARE the customer - we understand the pain deeply"
    ]

    personal_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(personal_y + 0.6), Inches(7), Inches(1.6)
    )
    ptf = personal_box.text_frame
    ptf.word_wrap = True

    for i, item in enumerate(personal_items):
        if i > 0:
            ptf.add_paragraph()
        p = ptf.paragraphs[i]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = white
        p.space_after = Pt(10)

    # Market Insight
    market_y = 3.2
    market_header = slide.shapes.add_textbox(
        Inches(1.5), Inches(market_y), Inches(7), Inches(0.4)
    )
    mh_frame = market_header.text_frame
    mh_frame.text = "🎯 Market Insight"
    p = mh_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = white

    market_items = [
        "Hundreds of student organizations at each university",
        "Each org books multiple events per semester",
        "Venues want consistent student group business",
        "No existing platform connects these two groups"
    ]

    market_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(market_y + 0.6), Inches(7), Inches(1.4)
    )
    mtf = market_box.text_frame
    mtf.word_wrap = True

    for i, item in enumerate(market_items):
        if i > 0:
            mtf.add_paragraph()
        p = mtf.paragraphs[i]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = white
        p.space_after = Pt(8)

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """Why are we building this? Because we've lived it. We're fraternity brothers at Ohio State, and we've personally experienced the pain of spending hours calling bar after bar trying to find one with availability for our events. We understand this problem intimately because we ARE the customer. We have a unique perspective as both the people experiencing the pain and the technical skills to solve it. The market opportunity is significant - there are hundreds of student organizations at every major university, each booking multiple events per semester. Venues want this consistent business from student groups, but there's no centralized platform connecting them. That's what we're building."""

    return slide


def create_closing_slide(prs, bg_color, white, gray, accent):
    """Slide 6: Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Header
    header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.6)
    )
    header_frame = header_box.text_frame
    header_frame.text = "VenueLink"
    p = header_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = white

    # Team
    team_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9), Inches(0.4)
    )
    tf = team_box.text_frame
    tf.text = "Michael Gorman & Larry Stelzer"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = white

    team_detail = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.6), Inches(9), Inches(0.3)
    )
    td_frame = team_detail.text_frame
    td_frame.text = "gorman.270 & stelzer.69"
    p = td_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.color.rgb = gray

    # Next Steps
    steps_y = 2.3
    steps_header = slide.shapes.add_textbox(
        Inches(1.5), Inches(steps_y), Inches(7), Inches(0.4)
    )
    sh_frame = steps_header.text_frame
    sh_frame.text = "Next Steps"
    p = sh_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = white

    steps = [
        "Complete venue discovery and booking flow (2 weeks)",
        "Conduct validation interviews with student orgs & venues",
        "Onboard first 5 venues and 3 student orgs at OSU",
        "Beta launch at Ohio State (3 months)"
    ]

    steps_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(steps_y + 0.6), Inches(7), Inches(1.8)
    )
    stf = steps_box.text_frame
    stf.word_wrap = True

    for i, step in enumerate(steps):
        if i > 0:
            stf.add_paragraph()
        p = stf.paragraphs[i]
        p.text = f"• {step}"
        p.font.size = Pt(16)
        p.font.color.rgb = white
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.CENTER

    # Call to Action
    cta_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.6), Inches(7), Inches(0.6)
    )
    cta_frame = cta_box.text_frame
    cta_frame.text = "Let's connect if you're interested in early venue partnerships\nor want to be part of shaping this platform"
    p = cta_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = white

    # Add speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """We're Michael Gorman and Larry Stelzer, and we're building VenueLink. Our immediate next steps are to finish the core booking flow in the next two weeks, conduct validation interviews with student organizations and venues, onboard our first partners, and launch a beta at Ohio State within three months. We're looking for early venue partners who want to reach more student organizations and student org leaders who want to help shape this platform. If you're interested in learning more or getting involved, let's connect."""

    return slide


if __name__ == "__main__":
    create_venuelink_pitch()
